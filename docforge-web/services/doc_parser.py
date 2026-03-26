"""참고 양식 파일 파싱 — .md/.txt/.html/.docx/.pdf → 텍스트 추출."""

from __future__ import annotations

import base64
import logging
from pathlib import Path

from services import codebase_parser

logger = logging.getLogger(__name__)


def parse_reference(data_base64: str, filename: str) -> str:
    """base64 인코딩된 파일에서 텍스트 추출. 최대 30000자.

    filename이 디렉터리 경로이면 codebase_parser에 위임한다.
    """
    # 로컬 디렉터리 경로인 경우 코드베이스 파서로 위임
    candidate = Path(filename)
    if candidate.is_dir():
        return codebase_parser.parse_codebase(str(candidate))

    ext = Path(filename).suffix.lower()
    raw = base64.standard_b64decode(data_base64)

    if ext in (".md", ".txt", ".markdown"):
        text = raw.decode("utf-8", errors="replace")

    elif ext in (".html", ".htm"):
        text = _parse_html(raw)

    elif ext == ".docx":
        text = _parse_docx(raw)

    elif ext == ".pptx":
        text = _parse_pptx(raw)

    elif ext == ".pdf":
        text = _parse_pdf(raw)

    else:
        text = raw.decode("utf-8", errors="replace")

    return text[:30000]


def _parse_html(raw: bytes) -> str:
    """HTML에서 구조(제목·목록·표·이미지·코드)를 마크다운으로 변환."""
    from bs4 import BeautifulSoup

    if isinstance(raw, str):
        soup = BeautifulSoup(raw, "html.parser")
    else:
        soup = BeautifulSoup(raw, "html.parser")

    # 불필요한 태그 제거
    for tag in soup(["script", "style", "nav", "footer", "iframe", "noscript"]):
        tag.decompose()

    parts: list[str] = []

    def _walk(element):
        if isinstance(element, str):
            return
        if not hasattr(element, "name") or element.name is None:
            return

        tag = element.name

        # 제목
        if tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            level = int(tag[1])
            text = element.get_text(strip=True)
            if text:
                parts.append(f"{'#' * level} {text}")
            return

        # 이미지
        if tag == "img":
            alt = element.get("alt", "")
            parts.append(f"[이미지: {alt}]" if alt else "[이미지 위치]")
            return

        # 표
        if tag == "table":
            rows: list[list[str]] = []
            for tr in element.find_all("tr"):
                cells = []
                for td in tr.find_all(["td", "th"]):
                    cells.append(td.get_text(strip=True))
                if cells:
                    rows.append(cells)
            if rows:
                header = "| " + " | ".join(rows[0]) + " |"
                separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
                body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
                parts.append(f"{header}\n{separator}\n{body}")
            return

        # 목록
        if tag in ("ul", "ol"):
            for i, li in enumerate(element.find_all("li", recursive=False), 1):
                text = li.get_text(strip=True)
                if text:
                    prefix = f"{i}." if tag == "ol" else "-"
                    parts.append(f"{prefix} {text}")
            return

        # 코드 블록
        if tag == "pre":
            code = element.get_text()
            lang = ""
            code_el = element.find("code")
            if code_el:
                cls = code_el.get("class", [])
                for c in cls:
                    if c.startswith("language-"):
                        lang = c.replace("language-", "")
                        break
                code = code_el.get_text()
            parts.append(f"```{lang}\n{code.strip()}\n```")
            return

        # 인라인 코드
        if tag == "code" and element.parent and element.parent.name != "pre":
            parts.append(f"`{element.get_text()}`")
            return

        # blockquote
        if tag == "blockquote":
            text = element.get_text(strip=True)
            if text:
                parts.append(f"> {text}")
            return

        # 단락
        if tag == "p":
            text = element.get_text(strip=True)
            if text:
                parts.append(text)
            # 단락 안의 이미지
            for img in element.find_all("img"):
                alt = img.get("alt", "")
                parts.append(f"[이미지: {alt}]" if alt else "[이미지 위치]")
            return

        # 그 외 컨테이너는 자식 순회
        for child in element.children:
            _walk(child)

    # body가 있으면 body부터, 없으면 전체
    body = soup.find("body") or soup
    _walk(body)

    return "\n\n".join(parts)


def _parse_docx(raw: bytes) -> str:
    """docx에서 단락·표·이미지 위치를 순서대로 추출."""
    import io
    from docx import Document
    from docx.oxml.ns import qn

    doc = Document(io.BytesIO(raw))
    parts: list[str] = []

    for element in doc.element.body:
        tag = element.tag.split("}")[-1] if "}" in element.tag else element.tag

        if tag == "p":
            # 단락 처리
            # 스타일(제목) 감지
            style_el = element.find(qn("w:pPr"))
            heading = ""
            if style_el is not None:
                style_ref = style_el.find(qn("w:pStyle"))
                if style_ref is not None:
                    style_name = style_ref.get(qn("w:val"), "")
                    if "Heading" in style_name or "heading" in style_name:
                        level = "".join(c for c in style_name if c.isdigit()) or "2"
                        heading = "#" * int(level) + " "

            # 이미지 감지
            drawings = element.findall(".//" + qn("w:drawing"))
            images_in_para = len(drawings)

            text = element.text or ""
            # 모든 run의 텍스트 합치기
            for run in element.findall(".//" + qn("w:t")):
                if run.text:
                    text += run.text

            text = text.strip()
            if text:
                parts.append(f"{heading}{text}")
            if images_in_para:
                for _ in range(images_in_para):
                    parts.append("[이미지 위치]")

        elif tag == "tbl":
            # 표 처리
            rows = element.findall(qn("w:tr"))
            table_rows: list[list[str]] = []
            for row in rows:
                cells = []
                for cell in row.findall(qn("w:tc")):
                    cell_text = ""
                    for p in cell.findall(qn("w:p")):
                        for run in p.findall(".//" + qn("w:t")):
                            if run.text:
                                cell_text += run.text
                    cells.append(cell_text.strip())
                table_rows.append(cells)

            if table_rows:
                # 마크다운 표로 변환
                header = "| " + " | ".join(table_rows[0]) + " |"
                separator = "| " + " | ".join("---" for _ in table_rows[0]) + " |"
                body = "\n".join(
                    "| " + " | ".join(r) + " |" for r in table_rows[1:]
                )
                parts.append(f"{header}\n{separator}\n{body}")

    return "\n\n".join(parts)


def _parse_pptx(raw: bytes) -> str:
    """pptx에서 슬라이드별 텍스트·표·이미지 위치 추출."""
    import io
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation(io.BytesIO(raw))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = [f"## 슬라이드 {i}"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        # 제목 슬라이드의 큰 텍스트는 제목으로
                        if shape == slide.shapes[0] and i == 1:
                            slide_parts.append(f"# {text}")
                        else:
                            slide_parts.append(text)

            elif shape.has_table:
                table = shape.table
                rows: list[list[str]] = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append(cells)
                if rows:
                    header = "| " + " | ".join(rows[0]) + " |"
                    separator = "| " + " | ".join("---" for _ in rows[0]) + " |"
                    body = "\n".join("| " + " | ".join(r) + " |" for r in rows[1:])
                    slide_parts.append(f"{header}\n{separator}\n{body}")

            elif shape.shape_type == 13:  # Picture
                slide_parts.append("[이미지 위치]")

            elif shape.has_chart:
                slide_parts.append("[차트 위치]")

        parts.append("\n\n".join(slide_parts))

    return "\n\n---\n\n".join(parts)


def _parse_pdf(raw: bytes) -> str:
    import io
    from pypdf import PdfReader
    reader = PdfReader(io.BytesIO(raw))
    pages = []
    for page in reader.pages[:50]:  # 최대 50페이지
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n\n".join(pages)
